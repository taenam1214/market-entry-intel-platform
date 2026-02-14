from io import BytesIO
from datetime import datetime

try:
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib.colors import HexColor
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
    from reportlab.lib.enums import TA_LEFT, TA_CENTER
    HAS_REPORTLAB = True
except ImportError:
    HAS_REPORTLAB = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


def _get_score_color(score):
    """Return color based on score value (0-10)."""
    if score >= 7:
        return '#16a34a'  # green
    elif score >= 4:
        return '#ca8a04'  # yellow
    return '#dc2626'  # red


def _get_go_nogo(market_score, complexity_score):
    """Determine GO/NO-GO recommendation."""
    if market_score >= 7 and complexity_score <= 6:
        return 'GO', 'High'
    elif market_score >= 5 and complexity_score <= 7:
        return 'GO', 'Medium'
    elif market_score >= 4:
        return 'CONDITIONAL GO', 'Medium'
    return 'NO-GO', 'High'


def generate_executive_summary_pdf(report):
    """Generate executive summary PDF. Returns bytes."""
    if not HAS_REPORTLAB:
        raise ImportError("reportlab is required for PDF generation. Install with: pip install reportlab")

    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, topMargin=0.75*inch, bottomMargin=0.75*inch)
    styles = getSampleStyleSheet()

    title_style = ParagraphStyle('CustomTitle', parent=styles['Title'], fontSize=20, spaceAfter=6, textColor=HexColor('#111827'))
    heading_style = ParagraphStyle('CustomHeading', parent=styles['Heading2'], fontSize=14, spaceBefore=16, spaceAfter=8, textColor=HexColor('#111827'))
    body_style = ParagraphStyle('CustomBody', parent=styles['Normal'], fontSize=10, leading=14, textColor=HexColor('#374151'))
    label_style = ParagraphStyle('Label', parent=styles['Normal'], fontSize=9, textColor=HexColor('#6b7280'))

    story = []
    dashboard = report.dashboard_data or {}
    scores = report.detailed_scores or {}
    revenue = report.revenue_projections or {}
    insights = report.key_insights or []
    actions = report.recommended_actions or {}

    # Title
    story.append(Paragraph('EXECUTIVE SUMMARY', title_style))
    story.append(Paragraph(f'{report.company_name} — {report.target_market} Market Entry Analysis', body_style))
    story.append(Paragraph(f'Prepared: {datetime.now().strftime("%B %d, %Y")} | Confidential', label_style))
    story.append(Spacer(1, 20))

    # Company Info
    story.append(Paragraph('Company Profile', heading_style))
    company_data = [
        ['Company', report.company_name or 'N/A'],
        ['Industry', report.industry or 'N/A'],
        ['Target Market', report.target_market or 'N/A'],
        ['Company Size', report.company_size or 'N/A'],
        ['Annual Revenue', report.annual_revenue or 'N/A'],
    ]
    t = Table(company_data, colWidths=[2*inch, 4*inch])
    t.setStyle(TableStyle([
        ('TEXTCOLOR', (0, 0), (0, -1), HexColor('#6b7280')),
        ('TEXTCOLOR', (1, 0), (1, -1), HexColor('#111827')),
        ('FONTSIZE', (0, 0), (-1, -1), 10),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
        ('TOPPADDING', (0, 0), (-1, -1), 6),
    ]))
    story.append(t)
    story.append(Spacer(1, 16))

    # Key Metrics
    story.append(Paragraph('Key Performance Indicators', heading_style))
    market_score = dashboard.get('market_opportunity_score', 'N/A')
    complexity_score = dashboard.get('entry_complexity_score', 'N/A')
    metrics_data = [
        ['Market Opportunity Score', f'{market_score}/10'],
        ['Competitive Intensity', dashboard.get('competitive_intensity', 'N/A')],
        ['Entry Complexity Score', f'{complexity_score}/10'],
        ['Revenue Potential', dashboard.get('revenue_potential', 'N/A')],
        ['Y1 Revenue Target', revenue.get('year_1', 'N/A')],
        ['Y3 Revenue Target', revenue.get('year_3', 'N/A')],
    ]
    t = Table(metrics_data, colWidths=[3*inch, 3*inch])
    t.setStyle(TableStyle([
        ('TEXTCOLOR', (0, 0), (0, -1), HexColor('#6b7280')),
        ('TEXTCOLOR', (1, 0), (1, -1), HexColor('#111827')),
        ('FONTSIZE', (0, 0), (-1, -1), 10),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
        ('TOPPADDING', (0, 0), (-1, -1), 6),
    ]))
    story.append(t)
    story.append(Spacer(1, 16))

    # Key Insights
    if insights:
        story.append(Paragraph('Key Insights', heading_style))
        for insight in insights[:6]:
            title = insight.get('title', '')
            desc = insight.get('description', '')
            priority = insight.get('priority', 'medium')
            badge = f'[{priority.upper()}] ' if priority else ''
            story.append(Paragraph(f'<b>{badge}{title}</b>', body_style))
            story.append(Paragraph(desc, body_style))
            story.append(Spacer(1, 6))

    # Recommended Actions
    if actions:
        story.append(Paragraph('Recommended Actions', heading_style))
        for phase, label in [('immediate', 'Immediate (0-3 months)'), ('short_term', 'Short-Term (3-12 months)'), ('long_term', 'Long-Term (1-3 years)')]:
            items = actions.get(phase, [])
            if items:
                story.append(Paragraph(f'<b>{label}:</b>', body_style))
                if isinstance(items, list):
                    for item in items:
                        story.append(Paragraph(f'• {item}', body_style))
                else:
                    story.append(Paragraph(str(items), body_style))
                story.append(Spacer(1, 8))

    # Footer
    story.append(Spacer(1, 24))
    story.append(Paragraph('Prepared by KairosAI Market Entry Intelligence Platform', label_style))
    story.append(Paragraph('Confidential — Executive Use Only', label_style))

    doc.build(story)
    return buffer.getvalue()


def generate_go_nogo_pdf(report):
    """Generate GO/NO-GO decision PDF. Returns bytes."""
    if not HAS_REPORTLAB:
        raise ImportError("reportlab is required for PDF generation. Install with: pip install reportlab")

    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, topMargin=0.75*inch, bottomMargin=0.75*inch)
    styles = getSampleStyleSheet()

    title_style = ParagraphStyle('CustomTitle', parent=styles['Title'], fontSize=20, spaceAfter=6, textColor=HexColor('#111827'))
    heading_style = ParagraphStyle('CustomHeading', parent=styles['Heading2'], fontSize=14, spaceBefore=16, spaceAfter=8, textColor=HexColor('#111827'))
    body_style = ParagraphStyle('CustomBody', parent=styles['Normal'], fontSize=10, leading=14, textColor=HexColor('#374151'))
    label_style = ParagraphStyle('Label', parent=styles['Normal'], fontSize=9, textColor=HexColor('#6b7280'))

    story = []
    dashboard = report.dashboard_data or {}
    insights = report.key_insights or []

    market_score = int(dashboard.get('market_opportunity_score', 0) or 0)
    complexity_score = int(dashboard.get('entry_complexity_score', 0) or 0)
    decision, confidence = _get_go_nogo(market_score, complexity_score)
    decision_color = '#16a34a' if decision == 'GO' else '#ca8a04' if decision == 'CONDITIONAL GO' else '#dc2626'

    # Title
    story.append(Paragraph('GO / NO-GO DECISION FRAMEWORK', title_style))
    story.append(Paragraph(f'{report.company_name} — {report.target_market}', body_style))
    story.append(Paragraph(f'Prepared: {datetime.now().strftime("%B %d, %Y")}', label_style))
    story.append(Spacer(1, 20))

    # Decision Box
    decision_style = ParagraphStyle('Decision', parent=styles['Title'], fontSize=28, textColor=HexColor(decision_color), alignment=TA_CENTER)
    story.append(Paragraph(decision, decision_style))
    story.append(Paragraph(f'Confidence: {confidence}', ParagraphStyle('Conf', parent=label_style, alignment=TA_CENTER)))
    story.append(Spacer(1, 20))

    # Scores
    story.append(Paragraph('Assessment Scores', heading_style))
    scores_data = [
        ['Market Opportunity', f'{market_score}/10'],
        ['Entry Complexity', f'{complexity_score}/10'],
        ['Competitive Intensity', dashboard.get('competitive_intensity', 'N/A')],
        ['Revenue Potential', dashboard.get('revenue_potential', 'N/A')],
    ]
    t = Table(scores_data, colWidths=[3*inch, 3*inch])
    t.setStyle(TableStyle([
        ('TEXTCOLOR', (0, 0), (0, -1), HexColor('#6b7280')),
        ('FONTSIZE', (0, 0), (-1, -1), 10),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
        ('TOPPADDING', (0, 0), (-1, -1), 6),
    ]))
    story.append(t)
    story.append(Spacer(1, 16))

    # Risks
    risks = [i for i in insights if i.get('type') == 'risk']
    if risks:
        story.append(Paragraph('Risk Assessment', heading_style))
        for risk in risks:
            story.append(Paragraph(f'<b>{risk.get("title", "")}</b>: {risk.get("description", "")}', body_style))
            story.append(Spacer(1, 4))

    # Footer
    story.append(Spacer(1, 24))
    story.append(Paragraph('Prepared by KairosAI Market Entry Intelligence Platform', label_style))

    doc.build(story)
    return buffer.getvalue()


def generate_investment_memo_pdf(report):
    """Generate investment memo PDF. Returns bytes."""
    if not HAS_REPORTLAB:
        raise ImportError("reportlab is required for PDF generation. Install with: pip install reportlab")

    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, topMargin=0.75*inch, bottomMargin=0.75*inch)
    styles = getSampleStyleSheet()

    title_style = ParagraphStyle('CustomTitle', parent=styles['Title'], fontSize=20, spaceAfter=6, textColor=HexColor('#111827'))
    heading_style = ParagraphStyle('CustomHeading', parent=styles['Heading2'], fontSize=14, spaceBefore=16, spaceAfter=8, textColor=HexColor('#111827'))
    body_style = ParagraphStyle('CustomBody', parent=styles['Normal'], fontSize=10, leading=14, textColor=HexColor('#374151'))
    label_style = ParagraphStyle('Label', parent=styles['Normal'], fontSize=9, textColor=HexColor('#6b7280'))

    story = []
    dashboard = report.dashboard_data or {}
    revenue = report.revenue_projections or {}
    insights = report.key_insights or []
    actions = report.recommended_actions or {}
    market_score = int(dashboard.get('market_opportunity_score', 0) or 0)
    complexity_score = int(dashboard.get('entry_complexity_score', 0) or 0)

    # Title
    story.append(Paragraph('INVESTMENT MEMORANDUM', title_style))
    story.append(Paragraph(f'{report.company_name} — {report.target_market} Market Expansion', body_style))
    story.append(Paragraph(f'Prepared: {datetime.now().strftime("%B %d, %Y")} | Confidential', label_style))
    story.append(Spacer(1, 20))

    # Investment Thesis
    story.append(Paragraph('Investment Thesis', heading_style))
    story.append(Paragraph(
        f'{report.company_name} presents a {"compelling" if market_score >= 7 else "moderate" if market_score >= 5 else "limited"} '
        f'opportunity for strategic market expansion into {report.target_market}, within the {report.industry} sector.',
        body_style
    ))
    story.append(Spacer(1, 12))

    # Market Analysis
    story.append(Paragraph('Market Analysis', heading_style))
    market_data = [
        ['Market Opportunity Score', f'{market_score}/10'],
        ['Entry Complexity Score', f'{complexity_score}/10'],
        ['Competitive Intensity', dashboard.get('competitive_intensity', 'N/A')],
        ['Revenue Potential', dashboard.get('revenue_potential', 'N/A')],
    ]
    t = Table(market_data, colWidths=[3*inch, 3*inch])
    t.setStyle(TableStyle([
        ('TEXTCOLOR', (0, 0), (0, -1), HexColor('#6b7280')),
        ('FONTSIZE', (0, 0), (-1, -1), 10),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
    ]))
    story.append(t)
    story.append(Spacer(1, 16))

    # Financial Projections
    story.append(Paragraph('Financial Projections', heading_style))
    fin_data = [
        ['Y1 Revenue Target', revenue.get('year_1', 'TBD')],
        ['Y3 Revenue Target', revenue.get('year_3', 'TBD')],
        ['Y1 Market Share', revenue.get('market_share_y1', 'TBD')],
        ['Y3 Market Share', revenue.get('market_share_y3', 'TBD')],
        ['Expected ROI', '3-5x within 3 years' if market_score >= 7 else '2-3x within 3 years' if market_score >= 5 else '1-2x within 3 years'],
        ['Break-even', '12-18 months' if market_score >= 7 else '18-24 months' if market_score >= 5 else '24+ months'],
    ]
    t = Table(fin_data, colWidths=[3*inch, 3*inch])
    t.setStyle(TableStyle([
        ('TEXTCOLOR', (0, 0), (0, -1), HexColor('#6b7280')),
        ('FONTSIZE', (0, 0), (-1, -1), 10),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
    ]))
    story.append(t)
    story.append(Spacer(1, 16))

    # Risks
    risks = [i for i in insights if i.get('type') == 'risk']
    if risks:
        story.append(Paragraph('Risk Assessment', heading_style))
        for risk in risks:
            story.append(Paragraph(f'<b>{risk.get("title", "")}</b>: {risk.get("description", "")}', body_style))
            story.append(Spacer(1, 4))

    # Recommendations
    if actions:
        story.append(Paragraph('Strategic Recommendations', heading_style))
        for phase, label in [('immediate', 'Immediate'), ('short_term', 'Short-Term'), ('long_term', 'Long-Term')]:
            items = actions.get(phase, [])
            if items:
                story.append(Paragraph(f'<b>{label}:</b>', body_style))
                if isinstance(items, list):
                    for item in items:
                        story.append(Paragraph(f'• {item}', body_style))
                else:
                    story.append(Paragraph(str(items), body_style))
                story.append(Spacer(1, 6))

    # Footer
    story.append(Spacer(1, 24))
    story.append(Paragraph('Prepared by KairosAI Market Entry Intelligence Platform', label_style))
    story.append(Paragraph('Confidential — Executive Use Only', label_style))

    doc.build(story)
    return buffer.getvalue()


def generate_board_presentation_pptx(report):
    """Generate board presentation PPTX. Returns bytes."""
    if not HAS_PPTX:
        raise ImportError("python-pptx is required for PPTX generation. Install with: pip install python-pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    dashboard = report.dashboard_data or {}
    revenue = report.revenue_projections or {}
    insights = report.key_insights or []
    actions = report.recommended_actions or {}
    market_score = int(dashboard.get('market_opportunity_score', 0) or 0)
    complexity_score = int(dashboard.get('entry_complexity_score', 0) or 0)
    decision, confidence = _get_go_nogo(market_score, complexity_score)

    def add_title_slide(title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(17, 24, 39)
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(107, 114, 128)
        p2.alignment = PP_ALIGN.CENTER
        return slide

    def add_content_slide(title, bullet_points):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Title
        txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.5), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(17, 24, 39)
        # Content
        txBox2 = slide.shapes.add_textbox(Inches(0.75), Inches(1.75), Inches(11.5), Inches(5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for i, point in enumerate(bullet_points):
            p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            p.text = f'• {point}'
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(55, 65, 81)
            p.space_after = Pt(8)
        return slide

    # Slide 1: Title
    add_title_slide(
        f'{report.company_name}',
        f'{report.target_market} Market Entry Strategy\n{datetime.now().strftime("%B %Y")}'
    )

    # Slide 2: Executive Summary
    add_content_slide('Executive Summary', [
        f'Target Market: {report.target_market} ({report.industry})',
        f'Market Opportunity Score: {market_score}/10',
        f'Entry Complexity Score: {complexity_score}/10',
        f'Competitive Intensity: {dashboard.get("competitive_intensity", "N/A")}',
        f'Revenue Potential: {dashboard.get("revenue_potential", "N/A")}',
        f'Recommendation: {decision} (Confidence: {confidence})',
    ])

    # Slide 3: Financial Projections
    add_content_slide('Financial Projections', [
        f'Year 1 Revenue: {revenue.get("year_1", "TBD")}',
        f'Year 3 Revenue: {revenue.get("year_3", "TBD")}',
        f'Y1 Market Share Target: {revenue.get("market_share_y1", "TBD")}',
        f'Y3 Market Share Target: {revenue.get("market_share_y3", "TBD")}',
        f'Expected ROI: {"3-5x" if market_score >= 7 else "2-3x" if market_score >= 5 else "1-2x"} within 3 years',
        f'Break-even: {"12-18 months" if market_score >= 7 else "18-24 months" if market_score >= 5 else "24+ months"}',
    ])

    # Slide 4: Key Insights
    insight_bullets = []
    for ins in insights[:6]:
        priority = ins.get('priority', '')
        badge = f'[{priority.upper()}] ' if priority else ''
        insight_bullets.append(f'{badge}{ins.get("title", "")}: {ins.get("description", "")}')
    if insight_bullets:
        add_content_slide('Key Insights', insight_bullets)

    # Slide 5: Risks
    risk_bullets = [f'{r.get("title", "")}: {r.get("description", "")}' for r in insights if r.get('type') == 'risk']
    if risk_bullets:
        add_content_slide('Risk Assessment', risk_bullets[:6])

    # Slide 6: Recommended Actions
    action_bullets = []
    for phase, label in [('immediate', 'Immediate'), ('short_term', 'Short-Term'), ('long_term', 'Long-Term')]:
        items = actions.get(phase, [])
        if items:
            if isinstance(items, list):
                for item in items[:3]:
                    action_bullets.append(f'{label}: {item}')
            else:
                action_bullets.append(f'{label}: {items}')
    if action_bullets:
        add_content_slide('Strategic Recommendations', action_bullets[:8])

    # Slide 7: Next Steps
    add_content_slide('Next Steps', [
        'Board approval for market entry initiative',
        'Resource allocation and team assignment',
        'Detailed implementation planning',
        'Local market validation and partnership development',
        'Pilot program launch and monitoring',
    ])

    # Slide 8: Thank You
    add_title_slide('Thank You', 'Prepared by KairosAI Market Entry Intelligence Platform')

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
